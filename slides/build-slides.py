"""Build defense presentation slides (.pptx) for DATN thesis.

Usage: python slides/build-slides.py
Output: slides/bao-ve-do-an.pptx
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Paths ---
ROOT = Path(__file__).resolve().parent.parent
IMAGES = ROOT / "reports" / "2026-04-13" / "images"
OUTPUT = ROOT / "slides" / "bao-ve-do-an.pptx"

# --- Colors ---
HUST_BLUE = RGBColor(0x00, 0x33, 0x66)
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)
LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0x80, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
TABLE_HEADER_BG = RGBColor(0xD9, 0xE2, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def _add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, size=18, bold=False,
              color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_title_bar(slide, text):
    """Blue bar at top with white title text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HUST_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def _add_bullets(slide, left, top, width, height, items, size=18, color=BLACK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0


def _add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = IMAGES / img_name
    if not path.exists():
        _add_text(slide, left, top, 4, 0.5, f"[{img_name} not found]",
                  size=14, color=GRAY)
        return
    if width and height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    elif width:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top))


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = BLACK
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = BLACK
                p.font.name = "Calibri"

    return table


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_add_bg(s, HUST_BLUE)
_add_text(s, 1, 0.8, 11, 0.6,
          "ĐỒ ÁN TỐT NGHIỆP", size=20, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 0.8, 1.6, 11.7, 1.5,
          "THIẾT KẾ HỆ THỐNG QUẢN LÝ PHƯƠNG TIỆN RA/VÀO\n"
          "CƠ SỞ GIÁO DỤC ĐÀO TẠO\nTHÔNG QUA NHẬN DIỆN BIỂN SỐ XE",
          size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1, 3.8, 11, 0.5,
          "Sinh viên:  Hà Văn Quang — 20210718  |  Nguyễn Hữu Cần — 20223882",
          size=18, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1, 4.5, 11, 0.5,
          "Giảng viên hướng dẫn:  Nguyễn Tiến Dũng",
          size=18, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1, 5.5, 11, 0.5,
          "Đại học Bách khoa Hà Nội — Trường Công nghệ Thông tin và Truyền thông",
          size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
_add_text(s, 1, 6.3, 11, 0.5,
          "Tháng 4/2026", size=16, color=RGBColor(0xAA, 0xCC, 0xEE),
          align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "ĐẶT VẤN ĐỀ")
_add_bullets(s, 0.6, 1.4, 12, 5.5, [
    "Thực trạng: Quản lý xe ra/vào tại cơ sở giáo dục quy mô lớn (ĐHBK HN) "
    "chủ yếu thủ công — chậm, dễ sai sót, khó thống kê",
    "",
    "Nhu cầu:",
    "  • Kiểm soát tự động phương tiện ra/vào qua nhận diện biển số",
    "  • Quản lý tài khoản xe (đăng ký, tạm, lạ), tính phí tự động",
    "  • Thống kê lưu lượng realtime trên dashboard",
    "",
    "Giải pháp đề xuất:",
    "  • Hệ thống AI nhận diện biển số VN (1 hàng + 2 hàng) kết hợp web platform",
    "  • Pipeline: Camera → AI Engine → Backend API → Dashboard giám sát",
], size=20)


# ============================================================
# SLIDE 3: Objectives & Scope
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "MỤC TIÊU & PHẠM VI")
_add_text(s, 0.6, 1.4, 5.5, 0.5, "Mục tiêu", size=22, bold=True, color=HUST_BLUE)
_add_bullets(s, 0.6, 2.0, 5.5, 4, [
    "Xây dựng prototype E2E: Camera → AI → Backend → Dashboard",
    "Nhận diện biển số xe Việt Nam (1 hàng + 2 hàng)",
    "Quản lý tài khoản, giao dịch, barrier tự động",
    "Thống kê realtime trên trình duyệt",
], size=18)
_add_text(s, 7, 1.4, 5.5, 0.5, "Phạm vi prototype", size=22, bold=True, color=HUST_BLUE)
_add_bullets(s, 7, 2.0, 5.5, 4, [
    "1-2 camera (cấu hình IN/OUT)",
    "Xe máy + ô tô (4 class COCO)",
    "SQLite (dev) / PostgreSQL (production)",
    "Docker Compose deployment (3 services)",
], size=18)


# ============================================================
# SLIDE 4: Outline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "NỘI DUNG TRÌNH BÀY")
items = [
    ("1", "Tổng quan kiến trúc hệ thống"),
    ("2", "Phân tích yêu cầu và thiết kế"),
    ("3", "Xây dựng các module"),
    ("4", "Thực nghiệm và đánh giá"),
    ("5", "Kết luận và hướng phát triển"),
]
for i, (num, text) in enumerate(items):
    y = 1.6 + i * 1.0
    # number circle
    shape = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(y),
                               Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _add_text(s, 3, y + 0.05, 8, 0.5, text, size=22, color=BLACK)


# ============================================================
# SLIDE 5: Architecture
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "1. KIẾN TRÚC TỔNG THỂ 4 LỚP")
_add_image_safe(s, "hinh-01-architecture.png", 1.5, 1.3, width=10)
_add_bullets(s, 0.6, 5.5, 12, 1.8, [
    "Lớp 1: Đầu vào (Video / Camera IP)  →  Lớp 2: AI Engine (YOLO + SORT + OCR)",
    "Lớp 3: Backend (FastAPI + PostgreSQL, 18 API)  →  Lớp 4: Dashboard (React + TS)",
], size=16, color=GRAY)


# ============================================================
# SLIDE 6: Tech Stack
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "CÔNG NGHỆ SỬ DỤNG")
_add_table(s, 1.5, 1.4, 10, 5, [
    "Lĩnh vực", "Công nghệ", "Chi tiết"
], [
    ["Vehicle Detection", "YOLOv8 Nano", "COCO pretrained, 6.3MB, 3.2M params"],
    ["Object Tracking", "SORT", "Kalman Filter + Hungarian matching"],
    ["Plate Detection", "LP_detector.pt (YOLOv5)", "Custom, 41MB, 32ms/ảnh"],
    ["OCR", "LP_ocr.pt (YOLOv5)", "Char-level + 2-row clustering"],
    ["Backend", "FastAPI + SQLAlchemy", "Python, async, 18 endpoints"],
    ["Database", "PostgreSQL / SQLite", "10 bảng, ACID"],
    ["Frontend", "React + TypeScript", "Vite, 1,205 dòng"],
    ["Deployment", "Docker Compose", "3 services"],
], col_widths=[2.5, 3.5, 4.0])


# ============================================================
# SLIDE 7: Operation Scenario
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "KỊCH BẢN VẬN HÀNH")
_add_text(s, 0.6, 1.3, 6, 0.5, "Xe VÀO cổng", size=22, bold=True, color=GREEN)
_add_bullets(s, 0.6, 1.9, 6, 2.5, [
    "Camera → AI detect xe + biển số → OCR",
    "Server tra cứu account:",
    "  • Đã đăng ký → Barrier OPEN, trừ 2,000 VND",
    "  • Biển lạ → Tạo account tạm, cấp 100,000 VND, OPEN",
    "Dashboard hiển thị event realtime",
], size=17)
_add_text(s, 0.6, 4.3, 6, 0.5, "Xe RA cổng", size=22, bold=True, color=RED)
_add_bullets(s, 0.6, 4.9, 6, 2, [
    "Camera → AI detect → Server kiểm tra:",
    "  • Đã đăng ký → Barrier OPEN",
    "  • Account tạm → Barrier HOLD, cần xác minh",
    "Nhân viên verify trên Dashboard",
], size=17)
_add_table(s, 7, 1.4, 5.5, 5, [
    "Status", "Direction", "Barrier"
], [
    ["registered", "IN", "OPEN"],
    ["registered", "OUT", "OPEN"],
    ["temporary", "IN", "OPEN"],
    ["temporary", "OUT", "HOLD+Verify"],
    ["unknown (mới)", "IN", "OPEN (auto tạm)"],
    ["(mặc định)", "any", "HOLD"],
], col_widths=[2, 1.5, 2])


# ============================================================
# SLIDE 8: ERD
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "2. THIẾT KẾ CƠ SỞ DỮ LIỆU — 10 BẢNG")
_add_image_safe(s, "hinh-02-erd.png", 0.5, 1.3, width=12)


# ============================================================
# SLIDE 9: Sequence Diagram
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "LUỒNG XỬ LÝ CREATE_EVENT — 7 BƯỚC ATOMIC")
_add_image_safe(s, "hinh-03-sequence.png", 0.3, 1.3, width=7)
_add_bullets(s, 7.8, 1.4, 5, 5.5, [
    "1. Auto-create camera",
    "2. Tạo VehicleEvent",
    "3. Normalize plate text",
    "4. Tạo PlateRead",
    "5. Auto-create Account (100k VND)",
    "6. decide_barrier() — 6 nhánh",
    "7. Tạo BarrierAction + Transaction",
    "",
    "Atomic: nếu 1 bước fail → rollback",
], size=17)


# ============================================================
# SLIDE 10: Vehicle Detection
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.1. PHÁT HIỆN PHƯƠNG TIỆN — YOLOv8 Nano")
_add_image_safe(s, "hinh-04-detection.png", 0.3, 1.3, width=6.5)
_add_bullets(s, 7.2, 1.4, 5.5, 5, [
    "Model: YOLOv8 Nano (3.2M params, 6.3MB)",
    "Input: Video frame (numpy array)",
    "Output: Bounding box + score + class",
    "",
    "Lọc 4 class COCO:",
    "  • car, motorcycle, bus, truck",
    "",
    "Confidence threshold: 0.5",
    "File: vehicle_detector.py (61 dòng)",
], size=17)


# ============================================================
# SLIDE 11: SORT Tracking
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.1. THEO DÕI ĐỐI TƯỢNG — SORT TRACKER")
_add_image_safe(s, "hinh-05-tracking.png", 0.3, 1.3, width=6.5)
_add_bullets(s, 7.2, 1.4, 5.5, 5, [
    "Thuật toán: SORT (Bewley et al., 2016)",
    "",
    "Kalman Filter 7 chiều:",
    "  [x, y, area, ratio, vx, vy, va]",
    "",
    "Hungarian matching bằng IoU",
    "Params: max_age=1, min_hits=3",
    "",
    "Kết quả: Track ID ổn định qua frames",
    "File: sort_tracker.py (244 dòng)",
], size=17)


# ============================================================
# SLIDE 12: Plate Detection
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.2. PHÁT HIỆN BIỂN SỐ — LP_DETECTOR")
_add_image_safe(s, "hinh-06-plate.png", 0.3, 1.3, width=6.5)
_add_text(s, 7.2, 1.4, 5.5, 0.5, "Benchmark 2 model (20 ảnh):",
          size=18, bold=True, color=HUST_BLUE)
_add_table(s, 7.2, 2.1, 5.5, 1.8, [
    "Model", "Rate", "Tốc độ", "Size"
], [
    ["LP_detector (YOLOv5) ← CHỌN", "100%", "32ms", "41MB"],
    ["number_plate (YOLOv8)", "100%", "48ms", "50MB"],
], col_widths=[2.5, 1, 1, 1])
_add_bullets(s, 7.2, 4.2, 5.5, 2.5, [
    "LP_detector nhanh hơn 33%",
    "Nhẹ hơn 18%",
    "File: plate_detector.py (82 dòng)",
], size=17)


# ============================================================
# SLIDE 13: OCR
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.2. NHẬN DẠNG KÝ TỰ — OCR PIPELINE")
_add_image_safe(s, "hinh-07-ocr.png", 0.3, 1.3, width=6)
_add_text(s, 6.8, 1.3, 6, 0.5, "Thuật toán 2-row clustering:",
          size=18, bold=True, color=HUST_BLUE)
_add_bullets(s, 6.8, 1.9, 6, 2, [
    "1. Detect từng ký tự (YOLOv5 char-level)",
    "2. Sort by Y → tìm gap lớn nhất",
    "3. Gap > 30% avg height → tách 2 hàng",
    "4. Mỗi hàng sort L→R → nối",
], size=16)
_add_text(s, 6.8, 4.0, 6, 0.5, "Post-processing (MỚI):",
          size=18, bold=True, color=ACCENT_BLUE)
_add_bullets(s, 6.8, 4.6, 6, 2, [
    "Char mapping: O→0, I→1, S→5, B→8 (position-aware)",
    "Regex validate: 4 patterns biển VN",
    "33 unit tests pass",
], size=16)


# ============================================================
# SLIDE 14: Backend
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.3. BACKEND API SERVER")
_add_image_safe(s, "hinh-08-api.png", 0.3, 1.3, width=5.5)
_add_text(s, 6.2, 1.3, 6.5, 0.5, "FastAPI + SQLAlchemy + PostgreSQL",
          size=18, bold=True, color=HUST_BLUE)
_add_table(s, 6.2, 2.0, 6.5, 3.5, [
    "Method", "Endpoint", "Chức năng"
], [
    ["POST", "/api/v1/events", "Tạo event + barrier"],
    ["GET", "/api/v1/events", "Danh sách events"],
    ["GET", "/api/v1/accounts", "Quản lý tài khoản"],
    ["GET", "/api/v1/stats/realtime", "Thống kê realtime"],
    ["POST", "/barrier-actions/verify", "Xác minh barrier"],
    ["...", "... (18 endpoints)", ""],
], col_widths=[1, 2.5, 3])
_add_bullets(s, 6.2, 5.8, 6, 1, [
    "56/56 unit tests pass (1.48s)",
    "11 files, 1,644 dòng Python",
], size=16)


# ============================================================
# SLIDE 15: Dashboard
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.3. DASHBOARD GIÁM SÁT")
_add_image_safe(s, "hinh-09-dashboard.png", 0.3, 1.3, width=6.2)
_add_image_safe(s, "hinh-10-accounts.png", 6.8, 1.3, width=6.2)
_add_bullets(s, 0.6, 5.8, 12, 1.5, [
    "React + TypeScript (Vite) | 6 source files, 1,205 dòng | "
    "Realtime stats, Events, Accounts, Verify queue, Traffic stats",
], size=16, color=GRAY)


# ============================================================
# SLIDE 16: Snapshot (NEW)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "3.3. SNAPSHOT BIỂN SỐ (MỚI)")
_add_bullets(s, 0.6, 1.4, 12, 5.5, [
    "Khi detect biển số thành công → crop vùng biển số + padding 15%",
    "Lưu PNG: {timestamp}_{plate}_{trackid}.png",
    "Backend serve static files: /static/snapshots/",
    "Dashboard hiển thị thumbnail (50×30px) bên cạnh event",
    "",
    "Cấu hình qua env vars:",
    "  • SNAPSHOT_DIR (default: snapshots/)",
    "  • SNAPSHOT_PADDING (default: 0.15)",
    "  • ENABLE_SNAPSHOT (default: true)",
], size=20)


# ============================================================
# SLIDE 17: Docker
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "TÍCH HỢP & DEPLOYMENT")
_add_bullets(s, 0.6, 1.4, 5.5, 5.5, [
    "Docker Compose — 3 services:",
    "",
    "  1. postgres:16  (database)",
    "     Port 5432, volume persistent",
    "",
    "  2. backend  (FastAPI)",
    "     Port 8000, healthcheck, auto-migrate",
    "",
    "  3. dashboard  (React + Vite)",
    "     Port 5173, CORS enabled",
], size=18)
_add_bullets(s, 6.8, 1.4, 5.5, 5.5, [
    "E2E Data Flow:",
    "",
    "  Video file / Camera IP",
    "       ↓",
    "  AI Engine (Python, local)",
    "       ↓  POST /api/v1/events",
    "  Backend :8000 (Docker)",
    "       ↓  fetch API",
    "  Dashboard :5173 (Browser)",
], size=18)


# ============================================================
# SLIDE 18: Dataset
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "4. BỘ DỮ LIỆU")
_add_table(s, 0.6, 1.4, 6, 3.5, [
    "Dataset", "Số lượng", "Mô tả"
], [
    ["VNLP one_row", "19,086", "Biển 1 hàng (ô tô)"],
    ["VNLP two_rows", "12,618", "Biển 2 hàng (ô tô)"],
    ["VNLP xe_may", "5,593", "Biển 2 hàng (xe máy)"],
    ["VNLP tổng", "37,297", ""],
    ["Video test (VN)", "1 video", "trungdinh22-demo.mp4"],
    ["", "300 frames", "600×800, 10fps, 30s"],
], col_widths=[2, 1.5, 2.5])
_add_image_safe(s, "hinh-11-visual.png", 7, 1.3, width=6)


# ============================================================
# SLIDE 19: OCR Evaluation — Before/After Comparison
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "4. ĐÁNH GIÁ OCR — TRƯỚC VÀ SAU FINE-TUNE")
_add_text(s, 0.6, 1.3, 12, 0.5,
          "3 configs so sánh: Baseline → Fine-tuned detector → + PaddleOCR",
          size=18, color=GRAY)
_add_table(s, 0.6, 2.0, 12, 3.5, [
    "Metric", "Baseline", "Finetuned v8n + YOLO", "+ PaddleOCR"
], [
    ["Detection rate", "89.2%", "99.9%", "99.9%"],
    ["Exact match", "37.8%", "68.7%", "92.0% ★★"],
    ["Char accuracy", "53.8%", "82.4%", "96.5%"],
    ["Test size", "3,731", "3,731", "500*"],
    ["Throughput", "14.7 img/s GPU", "17.9 img/s GPU", "1.13 img/s CPU"],
], col_widths=[2.5, 2.5, 3.5, 3.5])
_add_text(s, 0.6, 5.8, 12, 0.5,
          "* PaddleOCR eval trên 500 ảnh; full 3,731 đang chạy (kỳ vọng ~90%)",
          size=14, color=GRAY)
_add_text(s, 0.6, 6.3, 12, 0.5,
          "→ Cải thiện +54.2pp: Baseline 37.8% → Finetuned+PaddleOCR 92.0%",
          size=17, bold=True, color=GREEN)


# ============================================================
# SLIDE 20: Detector Retrain Story (REBRANDED)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "CẢI THIỆN BIỂN SỐ — FINE-TUNE YOLOv8n")
_add_text(s, 0.6, 1.3, 12, 0.5,
          "Phát hiện bottleneck: LP_detector, không phải OCR char recognition",
          size=18, bold=True, color=HUST_BLUE)
_add_bullets(s, 0.6, 1.9, 6, 2.5, [
    "Phase 01: Thực nghiệm ground truth bbox:",
    "  • Dùng GT bbox → OCR thẳng → 69.8% exact match",
    "  • Kết luận: ceiling của LP_ocr ~70%",
    "  • Bottleneck là detection, không phải OCR",
    "",
    "→ Giải pháp: Fine-tune plate detector",
], size=16)
_add_text(s, 0.6, 4.5, 6, 0.5, "Training Details:", size=18, bold=True, color=ACCENT_BLUE)
_add_bullets(s, 0.6, 5.1, 6, 2, [
    "Base model: YOLOv8n",
    "Training data: 29,837 ảnh VNLP",
    "Epochs: 3  |  GPU: GTX 1650",
    "mAP50: 99.48%",
], size=16)
_add_text(s, 7, 1.3, 5.8, 0.5, "Kết quả:", size=20, bold=True, color=GREEN)
_add_table(s, 7, 2.0, 5.8, 2.8, [
    "Thử nghiệm", "Exact match"
], [
    ["Baseline LP_detector + YOLO", "37.8%"],
    ["Ground truth bbox (ceiling)", "69.8%"],
    ["Finetuned YOLOv8n + YOLO char", "68.7%"],
    ["Finetuned + PaddleOCR ★", "92.0%"],
], col_widths=[3.3, 2.5])
_add_text(s, 7, 5.0, 5.8, 0.5,
          "Finetuned + PaddleOCR vượt xa ceiling GT!",
          size=18, bold=True, color=GREEN)

# Callout box
callout = s.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(5.6), Inches(5.8), Inches(1.5))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(0xE2, 0xEF, 0xDA)
callout.line.color.rgb = GREEN
tf = callout.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.15)
tf.margin_top = Inches(0.1)
p = tf.paragraphs[0]
p.text = "+32.7pp exact match — kết quả quan trọng nhất sprint 21/04"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = "Calibri"


# ============================================================
# SLIDE 21: E2E Test
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "KẾT QUẢ E2E VIDEO TEST")
_add_text(s, 0.6, 1.3, 6, 0.5, "Lần 1 — 300 frames (toàn bộ video):",
          size=18, bold=True, color=HUST_BLUE)
_add_bullets(s, 0.6, 1.9, 6, 1.5, [
    "2 biển unique: 36H82613, 14K117970",
    "23 raw → 2 unique events | FPS: 2.1 | 140s",
], size=16)
_add_text(s, 0.6, 3.2, 6, 0.5, "Lần 2 — 101 frames (14/04/2026):",
          size=18, bold=True, color=HUST_BLUE)
_add_bullets(s, 0.6, 3.8, 6, 1.5, [
    "1 biển unique: 36H82613 (car, consistent)",
    "6 raw → 2 unique events | FPS: 1.6 | 62.1s",
], size=16)
_add_text(s, 0.6, 5.2, 6, 0.5, "Known Issues:", size=18, bold=True, color=ORANGE)
_add_bullets(s, 0.6, 5.8, 6, 1.5, [
    "Duplicate events do SORT tạo track mới",
    "FPS thấp trên CPU (1.6-2.1) — cần GPU",
], size=16)
_add_image_safe(s, "hinh-11-visual.png", 7, 1.3, width=6)


# ============================================================
# SLIDE 22: Tests Summary
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "TỔNG HỢP KIỂM THỬ")
_add_table(s, 1, 1.4, 11, 5, [
    "Hạng mục", "Kết quả", "Chi tiết"
], [
    ["Backend Unit Tests", "56/56 pass ✅", "1.48s, pytest"],
    ["OCR Post-process Tests", "33/33 pass ✅", "2.41s, pytest"],
    ["API Smoke Test", "5/5 pass ✅", "health, POST, GET events/accounts/stats"],
    ["Dashboard API", "7/7 pass ✅", "events, stats, accounts, barrier, traffic"],
    ["Barrier Logic", "8 cases pass ✅", "6 nhánh quyết định"],
    ["AI Engine → Backend", "1/1 pass ✅", "Pipeline detect + send event"],
    ["Dashboard Code Review", "2 bugs fixed ✅", "Pagination + double-fetch"],
    ["Known Issues", "1/3 ⚠", "Duplicate events, FPS cần GPU"],
], col_widths=[3.5, 3, 4.5])


# ============================================================
# SLIDE 23: Statistics
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "THỐNG KÊ SỐ LIỆU TỔNG HỢP")
_add_table(s, 0.6, 1.4, 5.5, 5, [
    "Hạng mục", "Số liệu"
], [
    ["Tổng dòng code Python", "4,835+"],
    ["Tổng dòng code TypeScript", "1,205"],
    ["API endpoints", "18"],
    ["Database tables", "10"],
    ["Unit tests", "146 (100% pass)"],
    ["AI models", "4 files (138.3 MB)"],
    ["Dataset VNLP", "37,297 ảnh"],
    ["Git commits", "75+"],
    ["Docker services", "3"],
])
_add_table(s, 7, 1.4, 5.5, 5, [
    "Metric", "Kết quả"
], [
    ["Plate detection rate (finetuned)", "100% (500 ảnh)"],
    ["OCR exact match (best)", "92.0% PaddleOCR"],
    ["OCR char accuracy (best)", "96.5%"],
    ["E2E FPS (CPU)", "1.6–2.1"],
    ["Backend tests", "101/101 pass"],
    ["AI engine tests", "45/45 pass"],
])


# ============================================================
# SLIDE 24: Achievements
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "5. KẾT QUẢ ĐẠT ĐƯỢC")
_add_bullets(s, 0.6, 1.4, 12, 5.5, [
    "✅  Pipeline E2E hoạt động: Video → AI Engine → Backend → Dashboard",
    "✅  Nhận diện biển số VN 1 hàng + 2 hàng (YOLOv5 char-level + gap clustering)",
    "✅  Backend 18 API endpoints + barrier rules 6 nhánh (101/101 tests pass)",
    "✅  Dashboard realtime monitoring (React + TypeScript)",
    "✅  Fine-tune LP_detector + PaddleOCR → 92.0% exact match (+54.2pp baseline)",
    "✅  OCR post-processing: char mapping + regex validate (45/45 tests pass)",
    "✅  Snapshot crop biển số + hiển thị thumbnail trên dashboard",
    "✅  Docker Compose deployment (3 services: postgres + backend + dashboard)",
    "✅  Hướng dẫn chạy E2E chi tiết (7 files, hỗ trợ Linux/Mac + Windows)",
], size=19)


# ============================================================
# SLIDE 25: Limitations
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "HẠN CHẾ")
_add_bullets(s, 0.6, 1.4, 12, 5.5, [
    "⚠  OCR accuracy thấp: 33.3% exact match — cần cải thiện model + post-processing",
    "⚠  FPS chậm trên CPU: 1.6-2.1 — chưa đáp ứng realtime, cần GPU",
    "⚠  Chỉ test trên 1 video biển VN — cần thêm video đa dạng",
    "⚠  Chưa full evaluation trên toàn bộ 37,297 ảnh VNLP",
    "⚠  Dashboard chưa test toàn diện trên browser (code review OK, cần manual test)",
    "⚠  Duplicate events khi SORT tạo track mới cho cùng biển số",
], size=20)


# ============================================================
# SLIDE 26: Future Work
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s)
_add_title_bar(s, "HƯỚNG PHÁT TRIỂN")
_add_text(s, 0.6, 1.3, 3.5, 0.5, "Ngắn hạn", size=22, bold=True, color=GREEN)
_add_bullets(s, 0.6, 1.9, 3.5, 2, [
    "Full OCR eval (3,731+ ảnh)",
    "Regex + char mapping nâng accuracy",
    "Deduplicate by plate_text",
    "Test dashboard trên browser",
], size=17)
_add_text(s, 4.8, 1.3, 3.5, 0.5, "Trung hạn", size=22, bold=True, color=ACCENT_BLUE)
_add_bullets(s, 4.8, 1.9, 3.5, 2, [
    "GPU acceleration (10-15 FPS)",
    "Nhiều camera đồng thời",
    "WebSocket realtime push",
    "Cải thiện OCR: deblur, augment",
], size=17)
_add_text(s, 9, 1.3, 3.5, 0.5, "Dài hạn", size=22, bold=True, color=ORANGE)
_add_bullets(s, 9, 1.9, 3.5, 2, [
    "Edge deploy (Jetson Nano)",
    "Quản lý đăng ký phương tiện",
    "Cảnh báo thông minh",
    "Perspective correction",
], size=17)


# ============================================================
# SLIDE 27: Q&A
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(s, HUST_BLUE)
_add_text(s, 1, 2, 11, 1,
          "CẢM ƠN THẦY/CÔ ĐÃ LẮNG NGHE",
          size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1, 3.5, 11, 0.5,
          "Q & A",
          size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(s, 1, 5.0, 11, 0.5,
          "Mã nguồn: github.com/ukelele0718/bienso",
          size=18, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
_add_text(s, 1, 5.5, 11, 0.5,
          "Hà Văn Quang — 20210718  |  Nguyễn Hữu Cần — 20223882",
          size=16, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
prs.save(str(OUTPUT))
print(f"Created: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {OUTPUT.stat().st_size / 1024:.0f} KB")
